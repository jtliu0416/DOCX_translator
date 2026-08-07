"""PPTX processing helpers.

Extract translatable text from slides and generate a bilingual presentation.
For each slide with translatable content, the original slide is preserved and a
translated duplicate is inserted directly after it. The implementation focuses
on visible slide content: text frames, placeholders, grouped shapes, and table
cells.
"""

import asyncio
import json
import logging
import posixpath
import re
import uuid
import zipfile
from collections import Counter
from collections.abc import Iterable
from copy import deepcopy
from xml.etree import ElementTree as ET
from typing import Any

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.opc.package import Part
from pptx.opc.packuri import PackURI
from pptx.parts.chart import ChartPart

PRESENTATION_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
OFFICE_DRAWING_2014_NS = "http://schemas.microsoft.com/office/drawing/2014/main"
POWERPOINT_2010_NS = "http://schemas.microsoft.com/office/powerpoint/2010/main"
RELATIONSHIP_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
CSLD_TAG = f"{{{PRESENTATION_NS}}}cSld"
SPTREE_TAG = f"{{{PRESENTATION_NS}}}spTree"
CUST_DATA_LST_TAG = f"{{{PRESENTATION_NS}}}custDataLst"
CREATION_ID_TAG = f"{{{OFFICE_DRAWING_2014_NS}}}creationId"
P14_CREATION_ID_TAG = f"{{{POWERPOINT_2010_NS}}}creationId"
P14_MOD_ID_TAG = f"{{{POWERPOINT_2010_NS}}}modId"
P14_UNIQUE_ID_TAGS = {P14_CREATION_ID_TAG, P14_MOD_ID_TAG}
PPT_SLIDE_PART_RE = re.compile(r"^ppt/slides/slide(\d+)\.xml$")
PARTNAME_SEQUENCE_RE = re.compile(r"^(.*?)(\d+)(\.[^./]+)$")

LOGGER = logging.getLogger(__name__)


def _contains_chinese(text: str) -> bool:
    return any("\u4e00" <= char <= "\u9fff" or "\u3400" <= char <= "\u4dbf" for char in text)


def _contains_english(text: str) -> bool:
    return any(("a" <= char <= "z") or ("A" <= char <= "Z") for char in text)


def _english_ratio(text: str) -> float:
    stripped = text.strip()
    if not stripped:
        return 0.0
    english_count = sum(1 for char in stripped if ("a" <= char <= "z") or ("A" <= char <= "Z"))
    return english_count / len(stripped)


def _analyze_translatability(text: str, source_lang: str) -> dict[str, bool]:
    has_chinese = _contains_chinese(text)
    has_english = _contains_english(text)

    if source_lang == "zh":
        has_source = has_chinese
    elif source_lang == "en":
        has_source = has_english
    else:
        has_source = bool(text.strip())

    return {
        "contains_chinese": has_chinese,
        "contains_english": has_english,
        "already_translated": False,
        "skip": not has_source,
    }


def _shape_at_path(shapes: Iterable[Any], shape_path: list[int]) -> Any:
    current_shapes = shapes
    shape = None
    for path_index, index in enumerate(shape_path):
        shape = current_shapes[index]
        if path_index < len(shape_path) - 1:
            current_shapes = shape.shapes
    return shape


def _iter_shapes(shapes: Iterable[Any], shape_path: list[int] | None = None):
    base_path = shape_path or []
    for index, shape in enumerate(shapes):
        current_path = [*base_path, index]
        yield shape, current_path
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _iter_shapes(shape.shapes, current_path)


def _paragraph_text(paragraph: Any) -> str:
    return "".join(run.text for run in paragraph.runs).strip()


def _extract_text_frame_units(
    text_frame: Any,
    unit_base: dict[str, Any],
    units: list[dict[str, Any]],
    index: int,
    source_lang: str,
    target_lang: str,
) -> int:
    for paragraph_index, paragraph in enumerate(text_frame.paragraphs):
        text = _paragraph_text(paragraph)
        if not text:
            continue

        language = _analyze_translatability(text, source_lang)
        units.append({
            "index": index,
            "type": unit_base["type"],
            "slide_index": unit_base["slide_index"],
            "shape_path": unit_base["shape_path"],
            "paragraph_index": paragraph_index,
            "text": text,
            "source_lang": source_lang,
            "target_lang": target_lang,
            **language,
            **{k: v for k, v in unit_base.items() if k not in {"type", "slide_index", "shape_path"}},
        })
        index += 1
    return index


def _extract_units_sync(
    pptx_path: str,
    output_json_path: str,
    source_lang: str,
    target_lang: str,
) -> dict:
    presentation = Presentation(pptx_path)
    units: list[dict[str, Any]] = []
    index = 0

    for slide_index, slide in enumerate(presentation.slides):
        for shape, shape_path in _iter_shapes(slide.shapes):
            if getattr(shape, "has_table", False):
                table = shape.table
                for row_index, row in enumerate(table.rows):
                    for column_index, cell in enumerate(row.cells):
                        index = _extract_text_frame_units(
                            cell.text_frame,
                            {
                                "type": "ppt_table_cell",
                                "slide_index": slide_index,
                                "shape_path": shape_path,
                                "row_index": row_index,
                                "column_index": column_index,
                            },
                            units,
                            index,
                            source_lang,
                            target_lang,
                        )
                continue

            if getattr(shape, "has_text_frame", False):
                index = _extract_text_frame_units(
                    shape.text_frame,
                    {
                        "type": "ppt_text",
                        "slide_index": slide_index,
                        "shape_path": shape_path,
                    },
                    units,
                    index,
                    source_lang,
                    target_lang,
                )

    result = {"units": units}
    with open(output_json_path, "w", encoding="utf-8") as handle:
        json.dump(result, handle, ensure_ascii=False, indent=2)
    return result


async def extract_slide_text(
    pptx_path: str,
    output_json_path: str,
    source_lang: str = "zh",
    target_lang: str = "en",
) -> dict:
    return await asyncio.to_thread(
        _extract_units_sync,
        pptx_path,
        output_json_path,
        source_lang,
        target_lang,
    )


def _load_translations(translations_json_path: str) -> dict[int, str]:
    with open(translations_json_path, "r", encoding="utf-8") as handle:
        payload = json.load(handle)

    translations = {}
    for item in payload.get("translations", []):
        translations[int(item["index"])] = item.get("text", "")
    return translations


def _load_units(units_json_path: str) -> list[dict[str, Any]]:
    with open(units_json_path, "r", encoding="utf-8") as handle:
        payload = json.load(handle)
    return payload.get("units", [])


def _replace_paragraph_text(paragraph: Any, translated_text: str) -> None:
    runs = list(paragraph.runs)
    if runs:
        runs[0].text = translated_text
        for run in runs[1:]:
            run.text = ""
        return

    paragraph.text = translated_text


def _unit_text_frame_on_slide(slide: Any, unit: dict[str, Any]) -> Any | None:
    shape = _shape_at_path(slide.shapes, unit["shape_path"])
    if unit.get("type") == "ppt_table_cell":
        table = shape.table
        return table.cell(unit["row_index"], unit["column_index"]).text_frame

    return shape.text_frame if getattr(shape, "has_text_frame", False) else None


def _unit_text_frame(presentation: Presentation, unit: dict[str, Any]) -> Any | None:
    slide_index = unit["slide_index"]
    if slide_index >= len(presentation.slides):
        return None

    return _unit_text_frame_on_slide(presentation.slides[slide_index], unit)


def _remap_relationship_ids(slide_element: Any, relationship_ids: dict[str, str]) -> None:
    if not relationship_ids:
        return

    relationship_attr_prefix = f"{{{RELATIONSHIP_NS}}}"
    for element in slide_element.iter():
        for attr_name, attr_value in list(element.attrib.items()):
            if attr_name.startswith(relationship_attr_prefix) and attr_value in relationship_ids:
                element.set(attr_name, relationship_ids[attr_value])


def _remove_custom_data(slide_element: Any) -> None:
    for element in list(slide_element.iter()):
        for child in list(element):
            if child.tag == CUST_DATA_LST_TAG:
                element.remove(child)


def _new_powerpoint_numeric_id(used_numeric_ids: set[str]) -> str:
    while True:
        value = str((uuid.uuid4().int % 4_294_967_295) + 1)
        if value not in used_numeric_ids:
            used_numeric_ids.add(value)
            return value


def _collect_powerpoint_numeric_ids(presentation: Presentation) -> set[str]:
    ids: set[str] = set()
    for slide in presentation.slides:
        for element in slide._element.iter():
            if element.tag in {P14_CREATION_ID_TAG, P14_MOD_ID_TAG} and "val" in element.attrib:
                ids.add(element.attrib["val"])
    return ids


def _refresh_creation_ids(slide_element: Any, used_numeric_ids: set[str]) -> None:
    # Cloned slides must not keep PowerPoint extension IDs from the source slide;
    # duplicate p14 IDs trigger PowerPoint's repair prompt on open.
    for element in slide_element.iter():
        if element.tag == CREATION_ID_TAG and "id" in element.attrib:
            element.set("id", "{" + str(uuid.uuid4()).upper() + "}")
        elif element.tag in P14_UNIQUE_ID_TAGS and "val" in element.attrib:
            element.set("val", _new_powerpoint_numeric_id(used_numeric_ids))


def _next_related_partname(package: Any, source_partname: Any) -> PackURI:
    partname = str(source_partname)
    match = PARTNAME_SEQUENCE_RE.match(partname)
    if match:
        template = f"{match.group(1)}%d{match.group(3)}"
    else:
        path, extension = posixpath.splitext(partname)
        template = f"{path}%d{extension}"

    return package.next_partname(template)


def _clone_generic_part(package: Any, source_part: Any) -> Part:
    return Part(
        _next_related_partname(package, source_part.partname),
        source_part.content_type,
        package,
        source_part.blob,
    )


def _clone_chart_part(package: Any, source_chart_part: ChartPart) -> ChartPart:
    cloned_chart_part = ChartPart.load(
        package.next_partname(ChartPart.partname_template),
        source_chart_part.content_type,
        package,
        source_chart_part.blob,
    )

    relationship_ids: dict[str, str] = {}
    for old_rid, relationship in source_chart_part.rels.items():
        if relationship.is_external:
            new_rid = cloned_chart_part.relate_to(
                relationship.target_ref,
                relationship.reltype,
                is_external=True,
            )
            relationship_ids[old_rid] = new_rid
            continue

        cloned_target = _clone_generic_part(package, relationship.target_part)
        new_rid = cloned_chart_part.relate_to(cloned_target, relationship.reltype)
        relationship_ids[old_rid] = new_rid

    _remap_relationship_ids(cloned_chart_part._element, relationship_ids)

    return cloned_chart_part


def _clone_relationship_target(package: Any, relationship: Any) -> Any:
    if isinstance(relationship.target_part, ChartPart):
        return _clone_chart_part(package, relationship.target_part)

    return relationship.target_part


def _relationship_source_part(rels_name: str) -> str | None:
    if rels_name == "_rels/.rels":
        return ""

    folder, rels_filename = posixpath.split(rels_name)
    if not folder.endswith("/_rels") or not rels_filename.endswith(".rels"):
        return None

    return posixpath.join(folder[: -len("/_rels")], rels_filename[: -len(".rels")])


def _resolve_relationship_target(source_part: str, target: str, target_mode: str | None) -> str | None:
    if target_mode == "External" or not target:
        return None

    if target.startswith("/"):
        return target.lstrip("/")

    return posixpath.normpath(posixpath.join(posixpath.dirname(source_part), target))


def _slide_part_sort_key(slide_part: str) -> int:
    match = PPT_SLIDE_PART_RE.match(slide_part)
    return int(match.group(1)) if match else 0


def _duplicate_powerpoint_unique_id_issues(package: zipfile.ZipFile, slide_parts: list[str]) -> list[str]:
    issues: list[str] = []
    seen: dict[tuple[str, str], str] = {}

    for slide_part in slide_parts:
        try:
            root = ET.fromstring(package.read(slide_part))
        except ET.ParseError:
            continue

        for element in root.iter():
            if element.tag not in P14_UNIQUE_ID_TAGS:
                continue

            value = element.attrib.get("val")
            if not value:
                continue

            key = (element.tag, value)
            if key in seen:
                tag_name = element.tag.rsplit("}", 1)[-1]
                issues.append(f"duplicate p14:{tag_name} value {value} in {seen[key]} and {slide_part}")
            else:
                seen[key] = slide_part

    return issues


def _validate_pptx_package(pptx_path: str) -> list[str]:
    issues: list[str] = []

    try:
        Presentation(pptx_path)
    except Exception as exc:
        issues.append(f"python-pptx cannot open package: {exc}")

    try:
        with zipfile.ZipFile(pptx_path) as package:
            part_names = package.namelist()
            part_name_set = set(part_names)

            duplicate_parts = [name for name, count in Counter(part_names).items() if count > 1]
            for part_name in duplicate_parts:
                issues.append(f"duplicate zip entry: {part_name}")

            relationship_roots: dict[str, ET.Element] = {}
            for part_name in part_names:
                if not (part_name.endswith(".xml") or part_name.endswith(".rels")):
                    continue

                try:
                    root = ET.fromstring(package.read(part_name))
                except ET.ParseError as exc:
                    issues.append(f"invalid XML in {part_name}: {exc}")
                    continue

                if part_name.endswith(".rels"):
                    relationship_roots[part_name] = root

            for rels_name, rels_root in relationship_roots.items():
                source_part = _relationship_source_part(rels_name)
                if source_part is None:
                    continue

                relationship_ids: set[str] = set()
                for relationship in rels_root:
                    relationship_id = relationship.attrib.get("Id")
                    if relationship_id:
                        if relationship_id in relationship_ids:
                            issues.append(f"duplicate relationship id {relationship_id} in {rels_name}")
                        relationship_ids.add(relationship_id)

                    resolved_target = _resolve_relationship_target(
                        source_part,
                        relationship.attrib.get("Target", ""),
                        relationship.attrib.get("TargetMode"),
                    )
                    if resolved_target and resolved_target not in part_name_set:
                        issues.append(
                            f"missing relationship target from {rels_name}: "
                            f"{relationship_id or '<no id>'} -> {resolved_target}"
                        )

            slide_parts = sorted(
                [part_name for part_name in part_names if PPT_SLIDE_PART_RE.match(part_name)],
                key=_slide_part_sort_key,
            )
            issues.extend(_duplicate_powerpoint_unique_id_issues(package, slide_parts))
    except (OSError, zipfile.BadZipFile) as exc:
        issues.append(f"cannot read PPTX zip package: {exc}")

    return issues


def _duplicate_slide(
    presentation: Presentation,
    source_slide: Any,
    used_numeric_ids: set[str],
) -> Any:
    duplicate_slide = presentation.slides.add_slide(source_slide.slide_layout)
    package = presentation.part.package

    relationship_ids: dict[str, str] = {}
    for old_rid, relationship in source_slide.part.rels.items():
        if relationship.reltype in {RT.NOTES_SLIDE, RT.TAGS}:
            continue

        if relationship.is_external:
            new_rid = duplicate_slide.part.relate_to(
                relationship.target_ref,
                relationship.reltype,
                is_external=True,
            )
        else:
            target_part = _clone_relationship_target(package, relationship)
            new_rid = duplicate_slide.part.relate_to(
                target_part,
                relationship.reltype,
            )
        relationship_ids[old_rid] = new_rid

    source_element = source_slide._element
    duplicate_element = duplicate_slide._element
    source_c_sld = next(child for child in source_element if child.tag == CSLD_TAG)
    duplicate_c_sld = next(child for child in duplicate_element if child.tag == CSLD_TAG)
    duplicate_sp_tree = duplicate_slide.shapes._spTree
    source_sp_tree = source_slide.shapes._spTree

    for child in list(duplicate_sp_tree):
        duplicate_sp_tree.remove(child)
    for child in source_sp_tree:
        duplicate_sp_tree.append(deepcopy(child))

    duplicate_c_sld.attrib.clear()
    duplicate_c_sld.attrib.update(source_c_sld.attrib)
    for child in list(duplicate_c_sld):
        duplicate_c_sld.remove(child)
    for child in source_c_sld:
        duplicate_c_sld.append(duplicate_sp_tree if child.tag == SPTREE_TAG else deepcopy(child))

    for child in list(duplicate_element):
        duplicate_element.remove(child)
    for child in source_element:
        duplicate_element.append(duplicate_c_sld if child.tag == CSLD_TAG else deepcopy(child))

    _remap_relationship_ids(duplicate_element, relationship_ids)
    _remove_custom_data(duplicate_element)
    _refresh_creation_ids(duplicate_element, used_numeric_ids)
    return duplicate_slide


def _slide_id_element_for_slide(presentation: Presentation, slide: Any) -> Any:
    for slide_id in presentation.slides._sldIdLst.sldId_lst:
        if presentation.part.related_slide(slide_id.rId) is slide:
            return slide_id
    raise ValueError("Slide is not registered in presentation")


def _reorder_slides(presentation: Presentation, ordered_slide_ids: list[Any]) -> None:
    slide_id_list = presentation.slides._sldIdLst
    for slide_id in list(slide_id_list.sldId_lst):
        slide_id_list.remove(slide_id)
    for slide_id in ordered_slide_ids:
        slide_id_list.append(slide_id)


def _replace_slide_translations(
    slide: Any,
    units: list[dict[str, Any]],
    translations: dict[int, str],
) -> None:
    for unit in units:
        if unit.get("skip"):
            continue

        translated_text = translations.get(int(unit["index"]), "").strip()
        if not translated_text:
            continue

        text_frame = _unit_text_frame_on_slide(slide, unit)
        if text_frame is None:
            continue

        paragraph_index = unit["paragraph_index"]
        if paragraph_index >= len(text_frame.paragraphs):
            continue

        _replace_paragraph_text(text_frame.paragraphs[paragraph_index], translated_text)


def _insert_translations_sync(
    pptx_path: str,
    translations_json_path: str,
    output_path: str,
    units_json_path: str,
) -> None:
    translations = _load_translations(translations_json_path)
    units = _load_units(units_json_path)

    presentation = Presentation(pptx_path)
    original_slides = list(presentation.slides)
    original_slide_ids = list(presentation.slides._sldIdLst.sldId_lst)
    used_numeric_ids = _collect_powerpoint_numeric_ids(presentation)
    units_by_slide: dict[int, list[dict[str, Any]]] = {}
    for unit in units:
        if unit.get("skip"):
            continue
        units_by_slide.setdefault(int(unit["slide_index"]), []).append(unit)

    translated_slide_ids: dict[int, Any] = {}
    for slide_index, slide in enumerate(original_slides):
        slide_units = units_by_slide.get(slide_index)
        if not slide_units:
            continue

        translated_slide = _duplicate_slide(presentation, slide, used_numeric_ids)
        _replace_slide_translations(translated_slide, slide_units, translations)
        translated_slide_ids[slide_index] = _slide_id_element_for_slide(presentation, translated_slide)

    ordered_slide_ids: list[Any] = []
    for slide_index, slide_id in enumerate(original_slide_ids):
        ordered_slide_ids.append(slide_id)
        translated_slide_id = translated_slide_ids.get(slide_index)
        if translated_slide_id is not None:
            ordered_slide_ids.append(translated_slide_id)

    _reorder_slides(presentation, ordered_slide_ids)

    presentation.save(output_path)


async def insert_slide_translations(
    pptx_path: str,
    translations_json_path: str,
    output_path: str,
    units_json_path: str,
) -> None:
    await asyncio.to_thread(
        _insert_translations_sync,
        pptx_path,
        translations_json_path,
        output_path,
        units_json_path,
    )


async def validate_pptx(pptx_path: str) -> bool:
    issues = await asyncio.to_thread(_validate_pptx_package, pptx_path)
    if issues:
        preview = "; ".join(issues[:8])
        if len(issues) > 8:
            preview += f"; ... {len(issues) - 8} more issue(s)"
        LOGGER.warning("PPTX validation failed for %s: %s", pptx_path, preview)
        return False

    return True
