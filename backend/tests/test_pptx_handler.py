import asyncio
import json
import os
import re
import sys
import tempfile
import zipfile
from pathlib import Path
from xml.etree import ElementTree as ET
from unittest import TestCase

from pptx import Presentation
from pptx.util import Inches

sys.path.insert(0, str(Path(__file__).resolve().parents[1]))

from app.services.pptx_handler import (  # noqa: E402
    POWERPOINT_2010_NS,
    PRESENTATION_NS,
    _clone_chart_part,
    _extract_units_sync,
    _insert_translations_sync,
    validate_pptx,
)


SLIDE_PART_RE = re.compile(r"^ppt/slides/slide(\d+)\.xml$")
P14_UNIQUE_TAGS = {
    f"{{{POWERPOINT_2010_NS}}}creationId",
    f"{{{POWERPOINT_2010_NS}}}modId",
}


ET.register_namespace("p", PRESENTATION_NS)
ET.register_namespace("p14", POWERPOINT_2010_NS)


def _create_pptx(path: Path, slide_texts: list[str]) -> None:
    presentation = Presentation()
    blank_layout = presentation.slide_layouts[6]

    for text in slide_texts:
        slide = presentation.slides.add_slide(blank_layout)
        text_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
        text_box.text = text

    presentation.save(path)


def _inject_p14_ids(path: Path, ids_by_slide: dict[int, list[tuple[str, str]]]) -> None:
    temp_path = path.with_suffix(".tmp.pptx")

    with zipfile.ZipFile(path) as source, zipfile.ZipFile(temp_path, "w", zipfile.ZIP_DEFLATED) as target:
        for item in source.infolist():
            data = source.read(item.filename)
            match = SLIDE_PART_RE.match(item.filename)
            if match and int(match.group(1)) in ids_by_slide:
                root = ET.fromstring(data)
                common_slide = root.find(f"{{{PRESENTATION_NS}}}cSld")
                if common_slide is None:
                    raise AssertionError(f"missing p:cSld in {item.filename}")

                extension_list = common_slide.find(f"{{{PRESENTATION_NS}}}extLst")
                if extension_list is None:
                    extension_list = ET.SubElement(common_slide, f"{{{PRESENTATION_NS}}}extLst")

                for tag_name, value in ids_by_slide[int(match.group(1))]:
                    extension = ET.SubElement(
                        extension_list,
                        f"{{{PRESENTATION_NS}}}ext",
                        {"uri": "{BB962C8B-B14F-4D97-AF65-F5344CB8AC3E}"},
                    )
                    ET.SubElement(extension, f"{{{POWERPOINT_2010_NS}}}{tag_name}", {"val": value})

                data = ET.tostring(root, encoding="utf-8", xml_declaration=True)

            target.writestr(item, data)

    os.replace(temp_path, path)


def _p14_values(path: Path) -> dict[str, list[tuple[str, str]]]:
    values: dict[str, list[tuple[str, str]]] = {}

    with zipfile.ZipFile(path) as package:
        slide_parts = sorted(
            [name for name in package.namelist() if SLIDE_PART_RE.match(name)],
            key=lambda name: int(SLIDE_PART_RE.match(name).group(1)),
        )

        for slide_part in slide_parts:
            root = ET.fromstring(package.read(slide_part))
            for element in root.iter():
                if element.tag not in P14_UNIQUE_TAGS:
                    continue

                tag_name = element.tag.rsplit("}", 1)[-1]
                values.setdefault(tag_name, []).append((element.attrib["val"], slide_part))

    return values


class PptxHandlerTest(TestCase):
    def test_insert_translations_refreshes_cloned_powerpoint_ids(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            base = Path(temp_dir)
            original_path = base / "original.pptx"
            units_path = base / "units.json"
            translations_path = base / "translations.json"
            output_path = base / "translated.pptx"

            _create_pptx(original_path, ["项目进展"])
            _inject_p14_ids(
                original_path,
                {1: [("creationId", "12345"), ("modId", "67890")]},
            )

            units = _extract_units_sync(str(original_path), str(units_path), "zh", "en")["units"]
            with open(translations_path, "w", encoding="utf-8") as handle:
                json.dump(
                    {"translations": [{"index": unit["index"], "text": "Project progress"} for unit in units]},
                    handle,
                )

            _insert_translations_sync(str(original_path), str(translations_path), str(output_path), str(units_path))

            self.assertEqual(len(Presentation(output_path).slides), 2)
            self.assertTrue(asyncio.run(validate_pptx(str(output_path))))

            for occurrences in _p14_values(output_path).values():
                values = [value for value, _slide_part in occurrences]
                self.assertEqual(len(values), len(set(values)))

    def test_validate_pptx_rejects_duplicate_powerpoint_ids(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            path = Path(temp_dir) / "duplicate-p14.pptx"

            _create_pptx(path, ["第一页", "第二页"])
            _inject_p14_ids(
                path,
                {
                    1: [("creationId", "12345")],
                    2: [("creationId", "12345")],
                },
            )

            self.assertFalse(asyncio.run(validate_pptx(str(path))))

    def test_clone_chart_part_remaps_internal_relationship_ids(self) -> None:
        fixture = Path("backend/uploads/fc7f22b4-c253-4cfb-b915-648634044957/original.pptx")
        if not fixture.exists():
            self.skipTest("chart fixture is not available")

        presentation = Presentation(fixture)
        chart_relationship = next(
            relationship
            for relationship in presentation.slides[3].part.rels.values()
            if relationship.reltype.endswith("/chart")
        )

        cloned_chart = _clone_chart_part(presentation.part.package, chart_relationship.target_part)
        relationship_targets = {
            relationship.rId: relationship.target_part.content_type
            for relationship in cloned_chart.rels.values()
            if not relationship.is_external
        }

        root = ET.fromstring(cloned_chart.blob)
        external_data = next(
            element
            for element in root.iter()
            if element.tag.endswith("}externalData")
        )
        workbook_rid = external_data.attrib["{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"]

        self.assertEqual(
            relationship_targets[workbook_rid],
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        )
